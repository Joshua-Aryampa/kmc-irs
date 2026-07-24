"""Generate KMC IRS supervisor presentation (5 slides)."""
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

OUTPUT = Path(__file__).resolve().parents[1] / "docs" / "KMC-IRS-Presentation.pptx"
BRAND_RED = RGBColor(185, 28, 28)
DARK = RGBColor(15, 23, 42)
MUTED = RGBColor(100, 116, 139)
PLACEHOLDER = RGBColor(148, 163, 184)


def set_run(run, *, size=18, bold=False, color=DARK):
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = "Calibri"


def add_title(slide, text, *, top=0.35, size=32):
    box = slide.shapes.add_textbox(Inches(0.6), Inches(top), Inches(8.8), Inches(0.7))
    p = box.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = text
    set_run(run, size=size, bold=True, color=BRAND_RED)


def add_subtitle(slide, text, *, top=1.05):
    box = slide.shapes.add_textbox(Inches(0.6), Inches(top), Inches(8.8), Inches(0.45))
    p = box.text_frame.paragraphs[0]
    run = p.add_run()
    run.text = text
    set_run(run, size=16, color=MUTED)


def add_bullets(slide, items, *, left=0.75, top=1.55, width=5.0, height=4.5, size=16):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.level = 0
        p.space_after = Pt(8)
        run = p.add_run()
        run.text = item
        set_run(run, size=size, color=DARK)


def add_image_note(slide, text, *, left=6.0, top=1.55, width=3.3, height=4.8):
    shape = slide.shapes.add_shape(1, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(248, 250, 252)
    shape.line.color.rgb = RGBColor(226, 232, 240)
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    set_run(run, size=11, color=PLACEHOLDER)


def build():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    # Slide 1 — Title
    s1 = prs.slides.add_slide(blank)
    add_title(s1, "KMC Incident Reporting System", top=2.0, size=36)
    add_subtitle(s1, "Kiira Motors Corporation · Production IT", top=2.75)
    add_bullets(
        s1,
        [
            "Plant-wide web app for logging and tracking incidents",
            "Replaces paper-based reporting with a controlled digital workflow",
            "Built for shop-floor staff, supervisors, and management",
        ],
        top=3.45,
        width=8.5,
        size=18,
    )
    add_image_note(
        s1,
        "IMAGE: KMC logo\n(from login page or navbar)",
        left=7.0,
        top=0.55,
        width=2.4,
        height=1.2,
    )

    # Slide 2 — What was built
    s2 = prs.slides.add_slide(blank)
    add_title(s2, "What Was Built")
    add_bullets(
        s2,
        [
            "Secure login via Keycloak SSO (local demo mode for development)",
            "Dashboard — open/closed counts + recent incidents",
            "New incident — guided multi-step form with photos & witnesses",
            "Incident actions — verifier/approver queue with tabbed filters",
            "History — search, filter, paginate, CSV export",
            "Incident detail — full report, actions, signatories, timeline",
            "Responsive UI — works on desktop and mobile browsers",
        ],
        width=5.2,
        size=15,
    )
    add_image_note(
        s2,
        "IMAGE: Dashboard screenshot\n(stat cards + recent incidents table)",
        left=6.0,
        top=1.55,
    )

    # Slide 3 — Workflow
    s3 = prs.slides.add_slide(blank)
    add_title(s3, "Incident Workflow")
    add_bullets(
        s3,
        [
            "Draft → Pending verification → Pending approval → Closed",
            "Reporter submits → selects verifier → gets incident ID (MMYY###)",
            "Verifier sets severity, selects approver, verifies or rejects",
            "Approver closes case or rejects back to verifier",
            "Returned reports go back for correction and resubmit",
            "Email sent at submit, verify, reject, forward, and close",
            "Rule: reporter, verifier, and approver must be 3 different people",
        ],
        width=5.2,
        size=15,
    )
    add_image_note(
        s3,
        "IMAGE: Incident detail page\n(Pending verification — Report + Actions panels)\n\nOR simple workflow diagram",
        left=6.0,
        top=1.55,
    )

    # Slide 4 — Key features
    s4 = prs.slides.add_slide(blank)
    add_title(s4, "Key Features")
    add_bullets(
        s4,
        [
            "Drafts — save, edit, delete (private to reporter only)",
            "Employee search — type-ahead verifier/approver from Keycloak",
            "1–10 photos per report; at least 1 witness required",
            "Late submission flag + reason if > 30 min after incident time",
            "Digital signatures pulled from Keycloak user profile",
            "Full audit timeline on every incident",
            "PDF export for closed cases (KMC branding, photos, signatures)",
            "Role access — Employee, CEO (all + submit), Admin (all, no submit)",
        ],
        width=5.2,
        size=14,
    )
    add_image_note(
        s4,
        "IMAGE: New incident form\n(stepper + form fields)\n\nOR History page with filters",
        left=6.0,
        top=1.55,
    )

    # Slide 5 — Tech & readiness
    s5 = prs.slides.add_slide(blank)
    add_title(s5, "Technology & Readiness")
    add_bullets(
        s5,
        [
            "Stack: Python 3.11+, Django 5, PostgreSQL, Keycloak OIDC",
            "PDF generation, CSV export, email notifications built in",
            "UI aligned to KMC design spec across all main screens",
            "Automated tests with 80%+ coverage target on core modules",
            "Configurable via environment (.env) — ready for plant deployment",
            "Docs: README, SRS/SDD in docs/, demo accounts for UAT",
        ],
        width=5.2,
        size=15,
    )
    add_image_note(
        s5,
        "IMAGE: Closed incident PDF sample\n\nOR Incident actions queue tabs",
        left=6.0,
        top=1.55,
    )

    OUTPUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUTPUT)
    print(OUTPUT)


if __name__ == "__main__":
    build()
