"""Build IRS intern presentation from KMC sample template."""

from copy import deepcopy

from pptx import Presentation
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

TEMPLATE = r"c:\Users\ARYAMPA JOSHUA\Downloads\Sample Template.pptx"
OUTPUT = r"c:\Clean\IRS\kmc-irs\docs\IRS_Intern_Presentation.pptx"

TITLE_SLIDE = {
    "line1": "INCIDENT REPORTING SYSTEM (IRS)",
    "line2": "Intern Progress Presentation",
    "line3": "Information Systems Division",
}

SLIDES = [
    {
        "title": "2. Introduction",
        "bullets": [
            "Web-based Incident Reporting System (IRS) for Kiira Motors Corporation",
            "Digitises production incident reports from submission through verification to approval",
            "Built as a Django application with Keycloak single sign-on for plant employees",
            "Goal: one system usable plant-wide, with a clear audit trail and accountable sign-offs",
            "This presentation summarises the problem, solution, challenges, progress, and next steps",
        ],
    },
    {
        "title": "3. The Problem",
        "bullets": [
            "Incident reporting was manual, slow, and difficult to track across the plant",
            "Fixed production locations did not match a plant-wide reporting need",
            "Verifier and approver assignments were rigid and tied to old role/location rules",
            "Limited visibility into open incidents, late submissions, and workflow bottlenecks",
            "Hard to produce consistent closed reports with signatures and audit history",
        ],
    },
    {
        "title": "4. The Solution",
        "bullets": [
            "Django web app - works on desktop and mobile browsers",
            "Keycloak login for all employees; roles synced automatically on sign-in",
            "Reporter types scene location and selects a verifier; verifier selects an approver",
            "Workflow: Draft → Pending verification → Pending approval → Closed",
            "Dashboard with filters, personal queue, email notifications, CSV export, and PDF reports",
            "Signatories captured at submit, verify, and approve (signature image or printed name)",
        ],
    },
    {
        "title": "5. The Challenges",
        "bullets": [
            "Integrating company Keycloak (realm, client, redirect URI, service account for search)",
            "Replacing automatic routing with flexible employee search and manual sign-off selection",
            "Server/client clock sync for JWT validation during login",
            "PDF generation limits (table layout) and signature hosting still to be finalised",
            "Balancing intern delivery timeline with production-ready security and deployment needs",
        ],
    },
    {
        "title": "6. The Progress",
        "bullets": [
            "End-to-end workflow implemented and tested locally with Keycloak SSO",
            "Free-text incident location, pagination, and updated signatories section on reports",
            "Employee search for verifier (at submit) and approver (at verify)",
            "Software Requirements Specification and Software Design Document prepared",
            "UI improvements: clearer forms, branded logos, role-based dashboard visibility",
            "Closed incidents export to PDF with official KMC branding",
        ],
    },
    {
        "title": "7. The Next Steps",
        "bullets": [
            "Configure signature file storage (SIGNATURE_BASE_URL) with IT",
            "Production deployment: PostgreSQL, HTTPS, and application server setup",
            "User acceptance testing with supervisors and representative plant users",
            "Short user guide and handover to Information Systems / support team",
            "Monitor feedback and prioritise analytics enhancements for Admin and CEO roles",
        ],
    },
]


def _find_title_shape(slide):
    for shape in slide.shapes:
        if shape.is_placeholder and shape.placeholder_format.type == 1:
            return shape
        if shape.has_text_frame and "TITLE" in (shape.text_frame.text or "").upper():
            return shape
    for shape in slide.shapes:
        if shape.has_text_frame and shape.name.startswith("PlaceHolder"):
            return shape
    return None


def _find_title_slide_text(slide):
    best = None
    best_len = 0
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        text = shape.text_frame.text.strip()
        if len(text) > best_len:
            best = shape
            best_len = len(text)
    return best


def _set_title(shape, title):
    tf = shape.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = title
    p.alignment = PP_ALIGN.LEFT
    if p.runs:
        p.runs[0].font.bold = True
        p.runs[0].font.size = Pt(24)
        p.runs[0].font.name = "Poppins"


def _add_bullets(slide, bullets):
    left = Inches(0.65)
    top = Inches(1.55)
    width = Inches(11.5)
    height = Inches(5.2)
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.clear()
    for i, bullet in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = bullet
        p.level = 0
        p.space_after = Pt(10)
        p.alignment = PP_ALIGN.LEFT
        if p.runs:
            p.runs[0].font.size = Pt(18)
            p.runs[0].font.name = "Calibri"


def _duplicate_slide(prs, index):
    source = prs.slides[index]
    blank_layout = prs.slide_layouts[3]
    new_slide = prs.slides.add_slide(blank_layout)
    for shape in source.shapes:
        if shape.is_placeholder:
            continue
        new_el = deepcopy(shape.element)
        new_slide.shapes._spTree.insert_element_before(new_el, "p:extLst")
    return new_slide


def build():
    prs = Presentation(TEMPLATE)

    title_shape = _find_title_slide_text(prs.slides[0])
    if title_shape:
        tf = title_shape.text_frame
        tf.clear()
        lines = [TITLE_SLIDE["line1"], TITLE_SLIDE["line2"], TITLE_SLIDE["line3"]]
        for i, line in enumerate(lines):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = line
            p.alignment = PP_ALIGN.CENTER
            if p.runs:
                p.runs[0].font.bold = True
                p.runs[0].font.size = Pt(36 if i == 0 else 28)
                p.runs[0].font.name = "Poppins"

    content_template_index = 1
    while len(prs.slides) < len(SLIDES) + 1:
        _duplicate_slide(prs, content_template_index)

    for i, data in enumerate(SLIDES):
        slide = prs.slides[i + 1]
        title_shape = _find_title_shape(slide)
        if title_shape:
            _set_title(title_shape, data["title"])
        _add_bullets(slide, data["bullets"])

    prs.save(OUTPUT)
    print(f"Saved: {OUTPUT}")


if __name__ == "__main__":
    build()
