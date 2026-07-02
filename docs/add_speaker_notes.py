"""Add speaker notes to the IRS intern presentation (slides unchanged)."""

from pptx import Presentation

PRESENTATION = r"c:\Clean\IRS\kmc-irs\docs\IRS_Intern_Presentation.pptx"

SPEAKER_NOTES = [
    """Good morning/afternoon. We are Joshua Aryampa and Swahiba Nabisere, and this is our progress report on the Incident Reporting System (IRS).

In the next few minutes we will walk through why this system was needed, what we built, the main challenges we hit, what is working today, and what we need to do before it can go live plant-wide.""",

    """IRS is a web-based incident reporting tool for Kiira Motors — from first report through verification and final approval.

We chose Django for a maintainable backend and Keycloak so every employee signs in with the same company identity they already use. The real goal is not just digitizing a form: it is one plant-wide system with a clear audit trail and accountable sign-offs at each step.

This slide sets the roadmap for the rest of the presentation.""",

    """Today, manual or fragmented reporting makes incidents hard to track and easy to lose between people and departments.

Supervisors struggle to see what is still open, which submissions came in late, and where the workflow is stuck. When a case closes, there is often no single, consistent record with signatures and history — that weakens follow-up and accountability.

This is the gap IRS is meant to close.""",

    """Our solution is a browser-based app — usable on desktop or mobile without installing anything.

Employees log in through Keycloak; roles sync on sign-in. The reporter enters the scene location and chooses who verifies; the verifier then chooses the approver. That keeps responsibility explicit instead of hidden in fixed rules.

The workflow is Draft → Pending verification → Pending approval → Closed. The dashboard, queue, notifications, CSV export, and branded PDF give users one place to work and a proper closed report with signatories.""",

    """Most of our hardest work was integration, not layout. Keycloak needed the correct realm, client, redirect URI, and a service account so we could search employees when picking verifiers and approvers.

We also moved away from automatic routing to flexible search — a better fit for plant-wide reporting. JWT login failed until server and client clocks were aligned; that is a common SSO lesson, not a one-off bug.

Signature hosting is still being finalized with IT, which affects how signature images appear on closed reports.""",

    """We now have a working end-to-end flow tested locally with real Keycloak login — not a mock-up.

Key deliverables include free-text location, pagination, employee search for verifier and approver, an updated signatories section, role-based dashboard visibility, and PDF export with KMC branding. We also prepared the SRS and SDD so the work is documented for handover.

If helpful, we can briefly demo the dashboard or a closed PDF after the slides.""",

    """Near term, we are polishing the UI so the system feels ready for daily use.

Before production we need IT to configure signature storage (SIGNATURE_BASE_URL), deploy on PostgreSQL with HTTPS, and run user acceptance testing with supervisors and representative plant users. We will prepare a short user guide and hand over to Information Systems for support.

After launch, we will gather feedback and prioritize analytics improvements for Admin and CEO roles. We welcome your guidance on priorities and are happy to take questions.""",
]


def add_notes():
    prs = Presentation(PRESENTATION)
    if len(prs.slides) != len(SPEAKER_NOTES):
        raise SystemExit(
            f"Expected {len(SPEAKER_NOTES)} slides, found {len(prs.slides)}. "
            "Update SPEAKER_NOTES before running."
        )

    for slide, notes in zip(prs.slides, SPEAKER_NOTES):
        notes_frame = slide.notes_slide.notes_text_frame
        notes_frame.text = notes.strip()

    prs.save(PRESENTATION)
    print(f"Speaker notes added to {len(prs.slides)} slides: {PRESENTATION}")


if __name__ == "__main__":
    add_notes()
